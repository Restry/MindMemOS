"""文档抽取：把 Word / Excel / PPT / PDF / 文本 转成可入库的纯文本。

CLI（scripts/ingest/cli.py）和面板上传接口共用这一份实现，
避免两边解析行为不一致。

依赖：python-docx、openpyxl、python-pptx、pypdf（装在 MM 的 venv 里）。
老 .doc 走 macOS 自带的 textutil，不需要额外依赖。
"""

from __future__ import annotations

import os
import subprocess

MAX_CHARS = 12000  # 单条记忆上限，超了切片
MIN_CHARS = 80  # 太短的没有检索价值

EXTS = {".docx", ".doc", ".xlsx", ".xls", ".csv", ".pptx", ".ppt", ".pdf", ".md", ".txt", ".markdown"}


def from_docx(p: str) -> str:
    import docx

    d = docx.Document(p)
    parts = [x.text.strip() for x in d.paragraphs if x.text.strip()]
    # 表格也要抽——很多 Word 的关键信息都在表里，只抽段落会漏
    for t in d.tables:
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def from_doc(p: str) -> str:
    """老 .doc 用 macOS 自带 textutil，不用装额外东西。"""
    r = subprocess.run(["textutil", "-convert", "txt", "-stdout", p], capture_output=True, timeout=120)
    return r.stdout.decode("utf-8", "ignore")


def from_xlsx(p: str) -> str:
    """按行拼成「列名: 值」。

    为什么不整表塞：整张表丢进向量库，语义检索无从下手。拆成行之后
    "某设备的 IP 是多少" 这类问题才能精确命中到具体那一行。
    """
    import openpyxl

    wb = openpyxl.load_workbook(p, data_only=True, read_only=True)
    out = []
    for ws in wb.worksheets:
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        out.append(f"【工作表：{ws.title}】")
        header = [str(c).strip() if c is not None else "" for c in rows[0]]
        looks_header = sum(1 for h in header if h and not h.replace(".", "").isdigit()) >= 2
        for r in rows[1:] if looks_header else rows:
            vals = ["" if v is None else str(v).strip() for v in r]
            if not any(vals):
                continue
            if looks_header:
                pairs = [f"{h}: {v}" for h, v in zip(header, vals) if h and v]
                if pairs:
                    out.append("；".join(pairs))
            else:
                cells = [v for v in vals if v]
                if cells:
                    out.append(" | ".join(cells))
    wb.close()
    return "\n".join(out)


def from_csv(p: str) -> str:
    import csv

    out = []
    with open(p, newline="", encoding="utf-8", errors="ignore") as f:
        rows = list(csv.reader(f))
    if not rows:
        return ""
    header = [h.strip() for h in rows[0]]
    for r in rows[1:]:
        pairs = [f"{h}: {v.strip()}" for h, v in zip(header, r) if h and v.strip()]
        if pairs:
            out.append("；".join(pairs))
    return "\n".join(out)


def from_pptx(p: str) -> str:
    from pptx import Presentation

    prs = Presentation(p)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [sh.text_frame.text.strip() for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
        if texts:
            out.append(f"【第 {i} 页】" + " / ".join(texts))
    return "\n".join(out)


def from_pdf(p: str) -> str:
    from pypdf import PdfReader

    r = PdfReader(p)
    return "\n".join((pg.extract_text() or "").strip() for pg in r.pages)


def from_text(p: str) -> str:
    return open(p, encoding="utf-8", errors="ignore").read()


EXTRACT = {
    ".docx": from_docx,
    ".doc": from_doc,
    ".xlsx": from_xlsx,
    ".xls": from_xlsx,
    ".csv": from_csv,
    ".pptx": from_pptx,
    ".ppt": from_pptx,
    ".pdf": from_pdf,
    ".md": from_text,
    ".markdown": from_text,
    ".txt": from_text,
}


def extract(path: str) -> str:
    """按扩展名抽取纯文本。不支持的格式返回空串。"""
    fn = EXTRACT.get(os.path.splitext(path)[1].lower())
    return (fn(path) if fn else "") or ""


def chunks(text: str, limit: int = MAX_CHARS, min_chars: int = MIN_CHARS) -> list[str]:
    """按段落切片，单段超长再按字符硬切（否则整段会被丢掉）。"""
    out, cur = [], ""
    for para in text.split("\n"):
        if len(para) > limit:
            if cur.strip():
                out.append(cur.strip())
                cur = ""
            for i in range(0, len(para), limit):
                out.append(para[i : i + limit])
            continue
        if len(cur) + len(para) + 1 > limit:
            if cur.strip():
                out.append(cur.strip())
            cur = para
        else:
            cur += ("\n" if cur else "") + para
    if cur.strip():
        out.append(cur.strip())
    return [c for c in out if len(c) >= min_chars]
